from pptx.util import Inches
import math

class create_powerpoint():
    def __init__(self, prs, sentences, images) -> None:
        self.prs = prs
        self.ppt_text = sentences
        self.images = images

    def round_up(self, n, decimals=0):
        multiplier = 10 ** decimals
        return math.ceil(n * multiplier) / multiplier

    def presentation(self):
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        text_box = len(self.ppt_text)
        tex = 1.4
        height_to_begin_image = 1.5
        for i in range(text_box):
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(tex), Inches(9.5), Inches(0.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            points = tf.add_paragraph()
            points.text = self.ppt_text[i]
            
            if len(self.ppt_text[i]) < 13: #determining_gap_for_sentence
                points.level = 0
                tex += 0.6
                height_to_begin_image += 0.6
            elif len(self.ppt_text[i]) > 12: #determining_gap_for_sentence
                points.level = 0
                tex += 0.7
                height_to_begin_image += 0.7
            
        height_to_begin_image += 1
        items = list(self.images.items())
        if len(items) > 0 and len(items) < 5:
            shift = 0
            height_of_image = (7-height_to_begin_image)
            weidth_of_image = 9/(len(items))
            
            if height_of_image < weidth_of_image:
                weidth_of_image = height_of_image
            elif weidth_of_image < height_of_image:
                height_of_image = weidth_of_image
                
            for k in items:
                img = k[0]
                hyper_link = k[1]
                pic = slide.shapes.add_picture(img, Inches(0.25 + shift), Inches(height_to_begin_image), Inches(weidth_of_image), Inches(height_of_image))
                shift += weidth_of_image + 0.1
                pic.name = img
                
        elif len(items) > 4:        
            i_up = 0
            i_down = 0
            pos = 1
            height_of_image = ((7 - height_to_begin_image) / 2)
            images_in_line = self.round_up(len(items)/2)
            weidth_of_image = 9 / images_in_line
            
            if height_of_image < weidth_of_image:
                weidth_of_image = height_of_image
            elif weidth_of_image < height_of_image:
                height_of_image = weidth_of_image
                
            for k in items:
                img = k[0]
                hyper_link = k[1]
                if pos <= images_in_line:
                    pic = slide.shapes.add_picture(img, Inches(0.25 + i_up), Inches(height_to_begin_image), Inches(weidth_of_image), Inches(height_of_image))
                    pos += 1 
                    i_up += weidth_of_image + 0.1
                    pic.name = img

                elif pos > images_in_line:
                    pic = slide.shapes.add_picture(img, Inches(0.25 + i_down), Inches(height_to_begin_image + height_of_image+ 0.1), Inches(weidth_of_image), Inches(height_of_image))
                    pos += 1
                    i_down += weidth_of_image + 0.1
                    pic.name = img

        for shape in slide.shapes:
            for k in items:
                if shape.name == k[0]:
                    shape.click_action.hyperlink.address = k[1]
                    break
        return self.prs